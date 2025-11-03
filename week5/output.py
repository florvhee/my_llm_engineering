# In-Memory RAG System: No local file storage required!
# Everything processed in RAM: Download → Process → Vectorize → Delete

import os
import pickle
import io
import tempfile
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# LangChain imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import Chroma
from langchain.docstore.document import Document

# Processing libraries
import PyPDF2
import pandas as pd
from pptx import Presentation
import docx

# === AUTHENTICATION ===
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

def authenticate_drive(use_service_account=False):
    """Authenticate and return Google Drive service."""
    if use_service_account:
        from google.oauth2 import service_account
        creds = service_account.Credentials.from_service_account_file(
            'service-account-key.json', scopes=SCOPES)
        return build('drive', 'v3', credentials=creds)
    
    else:
        creds = None
        if os.path.exists('token.pickle'):
            with open('token.pickle', 'rb') as token:
                creds = pickle.load(token)
        
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    'credentials.json', SCOPES)
                creds = flow.run_local_server(port=0)
            
            with open('token.pickle', 'wb') as token:
                pickle.dump(creds, token)
        
        return build('drive', 'v3', credentials=creds)

# === IN-MEMORY FILE PROCESSORS ===

def process_pdf_from_bytes(file_bytes):
    """Extract text from PDF bytes without saving to disk."""
    try:
        pdf_file = io.BytesIO(file_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        return text
    except Exception as e:
        print(f"    ⚠️  Error processing PDF: {e}")
        return None

def process_docx_from_bytes(file_bytes):
    """Extract text from DOCX bytes without saving to disk."""
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = docx.Document(docx_file)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text
    except Exception as e:
        print(f"    ⚠️  Error processing DOCX: {e}")
        return None

def process_excel_from_bytes(file_bytes):
    """Extract text from Excel bytes without saving to disk."""
    try:
        excel_file = io.BytesIO(file_bytes)
        
        # Read all sheets
        xls = pd.ExcelFile(excel_file)
        text = ""
        
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            text += f"\n=== Sheet: {sheet_name} ===\n"
            text += df.to_string(index=False) + "\n"
        
        return text
    except Exception as e:
        print(f"    ⚠️  Error processing Excel: {e}")
        return None

def process_csv_from_bytes(file_bytes):
    """Extract text from CSV bytes without saving to disk."""
    try:
        csv_file = io.BytesIO(file_bytes)
        df = pd.read_csv(csv_file)
        return df.to_string(index=False)
    except Exception as e:
        print(f"    ⚠️  Error processing CSV: {e}")
        return None

def process_pptx_from_bytes(file_bytes):
    """Extract text from PowerPoint bytes without saving to disk."""
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n=== Slide {i} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text
    except Exception as e:
        print(f"    ⚠️  Error processing PowerPoint: {e}")
        return None

def process_text_from_bytes(file_bytes):
    """Extract text from plain text bytes."""
    try:
        # Try UTF-8 first
        return file_bytes.decode('utf-8')
    except:
        try:
            # Fallback to latin-1
            return file_bytes.decode('latin-1')
        except Exception as e:
            print(f"    ⚠️  Error processing text: {e}")
            return None

# === FILE PROCESSOR MAPPING ===
PROCESSORS = {
    'application/pdf': process_pdf_from_bytes,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': process_docx_from_bytes,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': process_excel_from_bytes,
    'application/vnd.ms-excel': process_excel_from_bytes,
    'text/csv': process_csv_from_bytes,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': process_pptx_from_bytes,
    'application/vnd.ms-powerpoint': process_pptx_from_bytes,
    'text/plain': process_text_from_bytes,
    'text/markdown': process_text_from_bytes,
    'text/html': process_text_from_bytes,
}

# === IN-MEMORY RAG SYSTEM ===
class InMemoryDriveRAG:
    def __init__(self, 
                 collection_name="drive_docs",
                 use_service_account=False,
                 chunk_size=1000,
                 chunk_overlap=200):
        """
        In-Memory RAG system - NO local file storage required!
        
        Files are downloaded to RAM, processed, vectorized, then discarded.
        Only the vector embeddings are persisted to disk.
        
        Args:
            collection_name: Name for the vector database collection
            use_service_account: Use service account auth instead of OAuth
            chunk_size: Size of text chunks for embeddings
            chunk_overlap: Overlap between chunks for context
        """
        self.service = authenticate_drive(use_service_account=use_service_account)
        
        # Use RecursiveCharacterTextSplitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )
        
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )
        
        self.vectorstore = None
        self.collection_name = collection_name
        
    def list_files(self, folder_id=None, file_types=None):
        """List files from Google Drive."""
        q = ["trashed=false"]
        
        if folder_id:
            q.append(f"'{folder_id}' in parents")
        
        if file_types:
            type_queries = [f"mimeType='{ft}'" for ft in file_types]
            q.append("(" + " or ".join(type_queries) + ")")
        
        query_string = " and ".join(q)
        
        results = self.service.files().list(
            q=query_string,
            pageSize=1000,
            fields="files(id, name, mimeType, size, modifiedTime)"
        ).execute()
        
        return results.get('files', [])
    
    def download_to_memory(self, file_id):
        """Download file directly to memory (BytesIO object)."""
        request = self.service.files().get_media(fileId=file_id)
        
        file_bytes = io.BytesIO()
        downloader = MediaIoBaseDownload(file_bytes, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_bytes.seek(0)  # Reset pointer to beginning
        return file_bytes.read()
    
    def export_to_memory(self, file_id, mime_type):
        """Export Google Docs/Sheets/Slides to memory."""
        export_map = {
            'application/vnd.google-apps.document': 'text/plain',
            'application/vnd.google-apps.spreadsheet': 'text/csv',
            'application/vnd.google-apps.presentation': 'text/plain',
        }
        
        export_mime = export_map.get(mime_type)
        if not export_mime:
            return None, None
        
        request = self.service.files().export_media(fileId=file_id, mimeType=export_mime)
        
        file_bytes = io.BytesIO()
        downloader = MediaIoBaseDownload(file_bytes, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_bytes.seek(0)
        
        # Return bytes and the actual MIME type for processing
        if export_mime == 'text/csv':
            return file_bytes.read(), 'text/csv'
        else:
            return file_bytes.read(), 'text/plain'
    
    def process_file_in_memory(self, file_info):
        """Download and process file entirely in memory."""
        file_id = file_info['id']
        file_name = file_info['name']
        mime_type = file_info['mimeType']
        
        print(f"  📥 Processing: {file_name}")
        
        try:
            # Handle Google Workspace files (export)
            if mime_type.startswith('application/vnd.google-apps'):
                file_bytes, actual_mime = self.export_to_memory(file_id, mime_type)
                if not file_bytes:
                    print(f"    ⚠️  Cannot export {mime_type}")
                    return None
                mime_type = actual_mime
            else:
                # Regular file download
                file_bytes = self.download_to_memory(file_id)
            
            # Process based on MIME type
            processor = PROCESSORS.get(mime_type)
            
            if processor:
                text = processor(file_bytes)
                
                if text:
                    return Document(
                        page_content=text,
                        metadata={
                            'source': file_name,
                            'file_id': file_id,
                            'mime_type': mime_type
                        }
                    )
            else:
                print(f"    ⚠️  No processor for {mime_type}")
                return None
                
        except Exception as e:
            print(f"    ❌ Error processing {file_name}: {e}")
            return None
    
    def build_vectorstore(self, folder_id=None, file_types=None, batch_size=10):
        """
        Stream-process files: Download → Process → Vectorize → Discard
        
        Args:
            folder_id: Google Drive folder ID (None for entire Drive)
            file_types: List of MIME types to include
            batch_size: Process files in batches to manage memory
        """
        print("🔍 Fetching file list from Google Drive...")
        files = self.list_files(folder_id=folder_id, file_types=file_types)
        print(f"📋 Found {len(files)} files to process\n")
        
        all_chunks = []
        processed_count = 0
        
        # Process files in batches
        for i in range(0, len(files), batch_size):
            batch = files[i:i + batch_size]
            print(f"📦 Processing batch {i//batch_size + 1}/{(len(files)-1)//batch_size + 1}")
            
            batch_documents = []
            
            # Process each file in the batch
            for file_info in batch:
                doc = self.process_file_in_memory(file_info)
                if doc:
                    batch_documents.append(doc)
                    processed_count += 1
            
            # Split into chunks
            if batch_documents:
                print(f"  ✂️  Splitting {len(batch_documents)} documents...")
                batch_chunks = self.text_splitter.split_documents(batch_documents)
                all_chunks.extend(batch_chunks)
                print(f"  ✅ Created {len(batch_chunks)} chunks from this batch\n")
            
            # Clear batch from memory
            batch_documents = None
        
        print(f"✅ Processed {processed_count}/{len(files)} files")
        print(f"✅ Created {len(all_chunks)} total chunks\n")
        
        # Create vector store (only this is persisted to disk)
        print("🔮 Creating vector database...")
        self.vectorstore = Chroma.from_documents(
            documents=all_chunks,
            embedding=self.embeddings,
            collection_name=self.collection_name,
            persist_directory="./chroma_db"
        )
        self.vectorstore.persist()
        print(f"✅ Vector database created!\n")
        
        return len(all_chunks)
    
    def query(self, question, k=4):
        """Query the RAG system."""
        if not self.vectorstore:
            try:
                self.vectorstore = Chroma(
                    collection_name=self.collection_name,
                    embedding_function=self.embeddings,
                    persist_directory="./chroma_db"
                )
            except:
                raise ValueError("No vector store found. Run build_vectorstore() first.")
        
        results = self.vectorstore.similarity_search(question, k=k)
        return results
    
    def query_with_scores(self, question, k=4):
        """Query with similarity scores."""
        if not self.vectorstore:
            self.vectorstore = Chroma(
                collection_name=self.collection_name,
                embedding_function=self.embeddings,
                persist_directory="./chroma_db"
            )
        
        results = self.vectorstore.similarity_search_with_score(question, k=k)
        return results

# === USAGE EXAMPLE ===
if __name__ == "__main__":
    # Initialize in-memory RAG system
    rag = InMemoryDriveRAG(
        chunk_size=1000,
        chunk_overlap=200
    )
    
    # Specify file types to process
    file_types = [
        # Documents
        'application/pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
        'application/vnd.google-apps.document',  # Google Docs
        
        # Spreadsheets
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # .xlsx
        'application/vnd.google-apps.spreadsheet',  # Google Sheets
        'text/csv',
        
        # Presentations
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # .pptx
        'application/vnd.google-apps.presentation',  # Google Slides
        
        # Text
        'text/plain',
        'text/markdown',
    ]
    
    # Process files (NO local storage!)
    # Option 1: Entire Drive
    # rag.build_vectorstore(file_types=file_types, batch_size=10)
    
    # Option 2: Specific folder
    folder_id = "YOUR_FOLDER_ID_HERE"
    rag.build_vectorstore(
        folder_id=folder_id, 
        file_types=file_types,
        batch_size=10  # Process 10 files at a time
    )
    
    # Query the system
    print("=" * 60)
    question = "What are the main topics in my documents?"
    results = rag.query_with_scores(question, k=3)
    
    print(f"❓ Question: '{question}'\n")
    for i, (doc, score) in enumerate(results, 1):
        print(f"📄 Result {i} (Similarity: {score:.4f})")
        print(f"   Source: {doc.metadata['source']}")
        print(f"   Content: {doc.page_content[:200]}...\n")